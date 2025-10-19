import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import docx
import pandas as pd
from bs4 import BeautifulSoup
from pptx import Presentation
import os
from sapperrag.read.base import BaseConversionStrategy


class ReadToolFacTory:
    def __init__(self):
        self.strategies = {
            "pdf": PDFToText(),
            "jpg": ImageToText(),
            "jpeg": ImageToText(),
            "png": ImageToText(),
            "bmp": ImageToText(),
            "gif": ImageToText(),
            "tiff": ImageToText(),
            "docx": DOCXToText(),
            "csv": CSVToText(),
            "xlsx": XLSXToText(),
            "md": MarkdownToText(),
            "html": HTMLToText(),
            "pptx": PPTXToText(),
            "txt": PlainText(),
        }

    def read_file(self, file_path):
        file_extension = os.path.splitext(file_path)[1].lower().strip('.')
        strategy = self.strategies.get(file_extension)
        if not strategy:
            raise ValueError(f"Unsupported file format: {file_extension}")
        return strategy.convert(file_path)


class PDFToText(BaseConversionStrategy):
    def convert(self, input_path):
        pdf_document = fitz.open(input_path)
        pdf_text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pdf_text += page.get_text()

        return pdf_text


class ImageToText(BaseConversionStrategy):
    def convert(self, input_path):
        image = Image.open(input_path)
        extracted_text = pytesseract.image_to_string(image)

        return extracted_text


class DOCXToText(BaseConversionStrategy):
    def convert(self, input_path):
        doc = docx.Document(input_path)
        doc_text = "\n".join([para.text for para in doc.paragraphs])

        return doc_text


class CSVToText(BaseConversionStrategy):
    def convert(self, input_path):
        df = pd.read_csv(input_path)
        formatted_rows = [
            ', '.join([f'{col}:{row[col]}' for col in df.columns])
            for _, row in df.iterrows()
        ]
        csv_text = '\n'.join(formatted_rows)

        return csv_text


class XLSXToText(BaseConversionStrategy):
    def convert(self, input_path):
        df = pd.read_excel(input_path)
        formatted_rows = [
            ', '.join([f'{col}:{row[col]}' for col in df.columns])
            for _, row in df.iterrows()
        ]
        xlsx_text = '\n\n'.join(formatted_rows)

        return xlsx_text


class MarkdownToText(BaseConversionStrategy):
    def convert(self, input_path):
        with open(input_path, "r", encoding="utf-8") as md_file:
            md_text = md_file.read()

        return md_text


class HTMLToText(BaseConversionStrategy):
    def convert(self, input_path):
        with open(input_path, "r", encoding="utf-8") as html_file:
            soup = BeautifulSoup(html_file, "html.parser")
        html_text = soup.get_text()

        return html_text


class PPTXToText(BaseConversionStrategy):
    def convert(self, input_path):
        prs = Presentation(input_path)
        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pptx_text += shape.text + "\n"

        return pptx_text


class PlainText(BaseConversionStrategy):
    def convert(self, input_path):
        with open(input_path, "r", encoding="utf-8") as txt_file:
            text = txt_file.read()

        return text
