import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import docx
import pandas as pd
from bs4 import BeautifulSoup
from pptx import Presentation
import os
from abc import ABC, abstractmethod
# from app import get_shared_state

# model_state = get_shared_state()


class FileConverter:
    def __init__(self):
        self.strategies = {
            "pdf": PDFToText(),
            "jpg": ImageToText(),
            "jpeg": ImageToText(),
            "png": ImageToText(),
            "bmp": ImageToText(),
            "gif": ImageToText(),
            "tiff": ImageToText(),
            "docx": DOCXToText(),
            "csv": CSVToText(),
            "xlsx": XLSXToText(),
            "md": MarkdownToText(),
            "html": HTMLToText(),
            "pptx": PPTXToText(),
            "txt": PlainText(),
        }

    def convert_to_text(self, file_path):
        file_extension = os.path.splitext(file_path)[1].lower().strip('.')
        strategy = self.strategies.get(file_extension)
        if not strategy:
            raise ValueError("Unsupported file format")
        return strategy.convert(file_path)


class ConversionStrategy(ABC):
    @abstractmethod
    def convert(self, input_path):
        pass


class PDFToText(ConversionStrategy):
    def convert(self, input_path):
        # model_list = model_state.model_list
        # pdf_text, images, out_meta = convert_single_pdf(input_path, model_list, 3)
        pdf_document = fitz.open(input_path)
        pdf_text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pdf_text += page.get_text()

        return pdf_text


class ImageToText(ConversionStrategy):
    def convert(self, input_path):
        image = Image.open(input_path)
        extracted_text = pytesseract.image_to_string(image)
        return extracted_text


class DOCXToText(ConversionStrategy):
    def convert(self, input_path):
        doc = docx.Document(input_path)
        doc_text = "\n".join([para.text for para in doc.paragraphs])
        return doc_text


class CSVToText(ConversionStrategy):
    def convert(self, input_path):
        df = pd.read_csv(input_path)
        # Initialize an empty list to store formatted strings
        formatted_rows = []
        for index, row in df.iterrows():
            # Format each row as col1:con1, col2:con2, col3:con3
            formatted_row = ', '.join([f'{col}:{row[col]}' for col in df.columns])
            formatted_rows.append(formatted_row)
        # Join all formatted rows into a single string with newline separators
        csv_text = '\n'.join(formatted_rows)
        return csv_text


class XLSXToText(ConversionStrategy):
    def convert(self, input_path):
        df = pd.read_excel(input_path)
        # Initialize an empty list to store formatted strings
        formatted_rows = []
        for index, row in df.iterrows():
            # Format each row as col1:con1, col2:con2, col3:con3
            formatted_row = ', '.join([f'{col}:{row[col]}' for col in df.columns])
            formatted_rows.append(formatted_row)
        # Join all formatted rows into a single string with newline separators
        xlsx_text = '\n\n'.join(formatted_rows)
        return xlsx_text


class MarkdownToText(ConversionStrategy):
    def convert(self, input_path):
        with open(input_path, "r", encoding="utf-8") as md_file:
            md_text = md_file.read()
        return md_text


class HTMLToText(ConversionStrategy):
    def convert(self, input_path):
        with open(input_path, "r", encoding="utf-8") as html_file:
            soup = BeautifulSoup(html_file, "html.parser")
        html_text = soup.get_text()
        return html_text


class PPTXToText(ConversionStrategy):
    def convert(self, input_path):
        prs = Presentation(input_path)
        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pptx_text += shape.text + "\n"
        return pptx_text


class PlainText(ConversionStrategy):
    def convert(self, input_path):
        with open(input_path, "r", encoding="utf-8") as txt_file:
            text = txt_file.read()
        return text


# 使用示例
# converter = FileConverter()
# converter.convert_to_text("example.pdf", "output.txt", "pdf")
# converter.convert_to_text("example.jpg", "output.txt", "jpg")
# converter.convert_to_text("example.docx", "output.txt", "docx")
