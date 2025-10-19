import aiohttp
import docx
import fitz
from aiohttp import ClientError
from urllib.parse import urlparse

from pptx import Presentation
from common.xfyun_tts import audio_to_text
from common.openai_image import image_to_text
from ..base_module.base_reader import PdfReader, DocxReader, PptxReader, XlsxReader, MdReader, CsvReader, TxtReader, \
    ImageReader, AudioReader
from .third_party_lib.docling.docling_server import TextExtractor

import os
import pandas as pd

# docling support
class PdfReaderBasedDocling(PdfReader):
    def __init__(self):
        super(PdfReaderBasedDocling, self).__init__()
        self.text_extractor = TextExtractor()

    async def read(self, file_path):
        return self.text_extractor.extract(file_path)


class DocxReaderBasedDocling(DocxReader):
    def __init__(self):
        super(DocxReaderBasedDocling, self).__init__()
        self.text_extractor = TextExtractor()

    async def read(self, file_path):
        return self.text_extractor.extract(file_path)


class PptxReaderBasedDocling(PptxReader):
    def __init__(self):
        super(PptxReaderBasedDocling, self).__init__()
        self.text_extractor = TextExtractor()

    async def read(self, file_path):
        return self.text_extractor.extract(file_path)


class XlsxReaderBasedDocling(XlsxReader):
    def __init__(self):
        super(XlsxReaderBasedDocling, self).__init__()
        self.text_extractor = TextExtractor()

    async def read(self, file_path):
        return self.text_extractor.extract(file_path)


class MdReaderBasedDocling(MdReader):
    def __init__(self):
        super(MdReaderBasedDocling, self).__init__()
        self.text_extractor = TextExtractor()

    async def read(self, file_path):
        return self.text_extractor.extract(file_path)


# python lib support
class CsvReaderBasedPythonLib(CsvReader):
    def __init__(self):
        super(CsvReaderBasedPythonLib, self).__init__()

    async def read(self, file_path):
        df = pd.read_csv(file_path)
        formatted_rows = [
            ', '.join([f'{col}:{row[col]}' for col in df.columns])
            for _, row in df.iterrows()
        ]
        return '\n'.join(formatted_rows)


class XlsxReaderBasedPythonLib(XlsxReader):
    def __init__(self):
        super(XlsxReaderBasedPythonLib, self).__init__()

    async def read(self, file_path):
        df = pd.read_excel(file_path)
        formatted_rows = [
            ', '.join([f'{col}:{row[col]}' for col in df.columns])
            for _, row in df.iterrows()
        ]
        return '\n'.join(formatted_rows)


class TxtReaderBasedPythonLib(TxtReader):
    def __init__(self):
        super(TxtReaderBasedPythonLib, self).__init__()
        self.session = None  # 我们将使用一个共享的aiohttp会话

    async def ensure_session(self):
        if self.session is None or self.session.closed:
            self.session = aiohttp.ClientSession()

    async def read(self, file_path):
        # 检查是否是URL（简单检查）
        if urlparse(file_path).scheme in ('http', 'https'):
            return await self._read_from_url(file_path)
        else:
            # 如果不是URL，保持原来的文件读取逻辑
            return await self._read_from_file(file_path)

    async def _read_from_url(self, url):
        await self.ensure_session()
        try:
            async with self.session.get(url) as response:
                response.raise_for_status()  # 如果响应状态不是200，抛出异常
                return await response.text(encoding='utf-8')
        except ClientError as e:
            raise IOError(f"Failed to read from URL {url}: {str(e)}")

    async def _read_from_file(self, file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as txt_file:
                return txt_file.read()
        except IOError as e:
            raise IOError(f"Failed to read from file {file_path}: {str(e)}")

    async def close(self):
        """关闭aiohttp会话"""
        if self.session and not self.session.closed:
            await self.session.close()


class PdfReaderBasedPythonLib(PdfReader):
    def __init__(self):
        super(PdfReaderBasedPythonLib, self).__init__()

    async def read(self, file_path):
        pdf_document = fitz.open(file_path)
        pdf_text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pdf_text += page.get_text()
        return pdf_text


class MdReaderBasedPythonLib(MdReader):
    def __init__(self):
        super(MdReaderBasedPythonLib, self).__init__()

    async def read(self, file_path):
        with open(file_path, "r", encoding="utf-8") as md_file:
            md_text = md_file.read()

        return md_text


class PptxReaderBasedPythonLib(PptxReader):
    def __init__(self):
        super(PptxReaderBasedPythonLib, self).__init__()

    async def read(self, file_path):
        prs = Presentation(file_path)
        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pptx_text += shape.text + "\n"

        return pptx_text


class DocxReaderBasedPythonLib(DocxReader):
    def __init__(self):
        super(DocxReaderBasedPythonLib, self).__init__()

    async def read(self, file_path):
        doc = docx.Document(file_path)
        doc_text = "\n".join([para.text for para in doc.paragraphs])

        return doc_text


# opeani support
class ImageReaderBasedOpenai(ImageReader):
    def __init__(self):
        super(ImageReaderBasedOpenai, self).__init__()

    async def read(self, file_path):
        content = await image_to_text(file_path)
        text = content.get('choices')
        if text:
            text = text[0].get('message').get('content')
        return text


class AudioReaderBasedXfyun(AudioReader):
    def __init__(self):
        super(AudioReaderBasedXfyun, self).__init__()

    async def read(self, file_path):
        content = await audio_to_text(file_path)
        return content


class FileReaderFactory:

    @staticmethod
    async def get_file_extension(file_path):
        print("file_path", file_path)
        return os.path.splitext(file_path)[-1].lstrip('.')

    @staticmethod
    async def get_file_reader(file_extension):
        if file_extension == "pdf":
            return PdfReaderBasedPythonLib()
        elif file_extension in {"docx", "doc"}:
            return DocxReaderBasedPythonLib()
        elif file_extension in {"pptx", "ppt"}:
            return PptxReaderBasedPythonLib()
        elif file_extension in {"xlsx", "xls"}:
            return XlsxReaderBasedPythonLib()
        elif file_extension == "md":
            return MdReaderBasedPythonLib()
        elif file_extension == "csv":
            return CsvReaderBasedPythonLib()
        elif file_extension == "txt":
            return TxtReaderBasedPythonLib()
        elif file_extension in {"png", "jpg"}:
            return ImageReaderBasedOpenai()
        elif file_extension in {"mp3","wav"}:
            return AudioReaderBasedXfyun()
        else:
            # Handle the case where the file extension is not recognized
            # raise ValueError(f"Unsupported file type: {file_extension}")
            print(f"Unsupported file type: {file_extension}")


